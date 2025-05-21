from pptx import Presentation

prs = Presentation('template.pptx')
slide_layout = prs.slide_layouts[1]  # Título e Conteúdo
slide = prs.slides.add_slide(slide_layout)

# Verifica se o título existe
if slide.shapes.title:
    slide.shapes.title.text = "Título de Exemplo"

# Verifica se o placeholder de conteúdo existe
if len(slide.placeholders) > 1:
    corpo = slide.placeholders[1]
    corpo.text = "Resumo de Exemplo\nhttps://exemplo.com\nData: 21/05/2025"

prs.save('teste_output.pptx')
