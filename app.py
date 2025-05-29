import re
from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import tempfile
import os

app = Flask(__name__)
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'template.pptx')

@app.route("/gerar_pptx", methods=["POST"])
def gerar_pptx():
    try:
        data = request.json
        prs = Presentation(TEMPLATE_PATH)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texto = shape.text

                    # Encontrar todos os {{campos}} presentes no texto
                    campos = re.findall(r"\{\{(\w+)\}\}", texto)
                    if campos:
                        tf = shape.text_frame
                        tf.clear()

                        for campo in campos:
                            valor = str(data.get(campo, f"{{{{{campo}}}}}")).replace("\\n", "\n")

                            # 🗓️ CONVERSÃO DE DATA ISO → BR
                            if "data" in campo.lower():
                                match_iso = re.match(r"(\d{4}-\d{2}-\d{2})", valor)  # Ajustado para pegar só a parte da data
                                if match_iso:
                                    try:
                                        dt = datetime.strptime(match_iso.group(1), "%Y-%m-%d")
                                        valor = dt.strftime("%d/%m/%Y")  # Apenas a data
                                    except Exception as e:
                                        print(f"⚠️ Erro ao converter data '{valor}': {e}")


                            p = tf.add_paragraph()
                            run = p.add_run()
                            run.text = valor

                            # 🎨 ESTILO POR CAMPO
                            if "titulo" in campo.lower():
                                run.font.bold = True
                                run.font.size = Pt(15)
                                run.font.color.rgb = RGBColor(124, 124, 124)
                                run.font.name = "Poppins"
                                p.alignment = PP_ALIGN.JUSTIFY

                            elif "data" in campo.lower():
                                run.font.italic = True
                                run.font.size = Pt(10)
                                run.font.color.rgb = RGBColor(124, 124, 124)
                                run.font.name = "Poppins"

                            elif "resumo" in campo.lower():
                                run.font.size = Pt(13)
                                run.font.color.rgb = RGBColor(124, 124, 124)
                                run.font.name = "Poppins"
                                p.alignment = PP_ALIGN.JUSTIFY

                            elif "link" in campo.lower():
                                run.font.size = Pt(8)
                                run.font.underline = True
                                run.font.name = "Poppins"
                                run.hyperlink.address = valor
                                run.font.color.rgb = RGBColor(255, 0, 0) 

        # 🔽 Salvar como arquivo temporário
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)

        return send_file(
            temp_file.name,
            as_attachment=True,
            download_name="noticias_geradas.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        print("🔥 Erro interno:", str(e))
        return {"erro": str(e)}, 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
