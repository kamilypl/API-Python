from flask import Flask, request, send_file
from pptx import Presentation
import tempfile
import os

app = Flask(__name__)
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'template.pptx')

@app.route("/gerar_pptx", methods=["POST"])
def gerar_pptx():
    try:
        data = request.json
        print("🔹 Dados recebidos:", data)

        prs = Presentation(TEMPLATE_PATH)
        slide_layout = prs.slide_layouts[1]  # Título e Conteúdo
        slide = prs.slides.add_slide(slide_layout)

        # Preencher título (com verificação segura)
        #if slide.shapes.title is not None:
            #slide.shapes.title.text = data.get('titulo', '')!!!!!!

        # Preencher corpo (com verificação segura)
        if len(slide.placeholders) > 1:
            corpo = slide.placeholders[1]
            corpo.text = f"{data.get('titulo', '')}\n\nData:{data.get('data', '')}\n\n{data.get('resumo', '')}\n\nFonte: {data.get('link', '')}"

        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_file.name)

        return send_file(
            temp_file.name,
            as_attachment=True,
            download_name="noticia.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        print("🔥 Erro interno:", str(e))
        return {"erro": str(e)}, 500

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
